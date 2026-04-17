"""Generate GitHub Advanced Security slide deck (PPTX + HTML)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_TEAL = RGBColor(0x2C, 0x4A, 0x52)
MED_TEAL = RGBColor(0x3A, 0x6B, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT = RGBColor(0x55, 0x55, 0x66)
A_ORANGE = RGBColor(0xE8, 0x6C, 0x3A)
A_GREEN = RGBColor(0x2E, 0x7D, 0x32)
A_BLUE = RGBColor(0x29, 0xB6, 0xF6)
LIGHT_TEXT = RGBColor(0xE0, 0xF0, 0xF5)
SUB_COLOR = RGBColor(0xBB, 0xDD, 0xE5)


def hdr(slide, title, sub=None):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
    h.fill.solid()
    h.fill.fore_color.rgb = DARK_TEAL
    h.line.fill.background()
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.8), prs.slide_width, Inches(0.06))
    a.fill.solid()
    a.fill.fore_color.rgb = A_BLUE
    a.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    if sub:
        p2 = tb.text_frame.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(18)
        p2.font.color.rgb = SUB_COLOR
        p2.font.name = "Segoe UI"


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_bullets(text_frame, bullets, font_size=14, color=None, prefix="\u2713  "):
    if color is None:
        color = LIGHT_TEXT
    text_frame.word_wrap = True
    for j, b in enumerate(bullets):
        bp = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
        bp.text = prefix + b
        bp.font.size = Pt(font_size)
        bp.font.color.rgb = color
        bp.font.name = "Segoe UI"
        bp.space_after = Pt(10)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_TEAL
bg.line.fill.background()

tb = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "GitHub Advanced Security"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Inches(0.06))
a.fill.solid()
a.fill.fore_color.rgb = A_BLUE
a.line.fill.background()

tb2 = s.shapes.add_textbox(Inches(2), Inches(3.7), Inches(9.3), Inches(1.0))
p2 = tb2.text_frame.paragraphs[0]
p2.text = "Securing Your Code at Every Stage of Development"
p2.font.size = Pt(24)
p2.font.color.rgb = SUB_COLOR
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

for i, lbl in enumerate(["Code Scanning", "Secret Scanning", "Dependency Review"]):
    x = Inches(2.5) + Inches(i * 3.2)
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(2.6), Inches(0.9))
    bx.fill.solid()
    bx.fill.fore_color.rgb = MED_TEAL
    bx.line.fill.background()
    tb3 = s.shapes.add_textbox(x + Inches(0.1), Inches(5.35), Inches(2.4), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = lbl
    p3.font.size = Pt(18)
    p3.font.color.rgb = WHITE
    p3.font.name = "Segoe UI"
    p3.font.bold = True
    p3.alignment = PP_ALIGN.CENTER

add_notes(s, "Welcome everyone. Today we explore GitHub Advanced Security \u2014 a developer-first suite of security tools that helps you find and fix vulnerabilities directly in your development workflow.")
print("Slide 1: Title")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — What is GHAS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "What is GitHub Advanced Security?")

features = [
    ("Developer-First Security", ["Security tools designed for developers", "No context switching required", "Shift-left approach to security"], DARK_TEAL),
    ("Automated Detection", ["Scans for vulnerabilities continuously", "Identifies secrets, code flaws & risky deps", "Runs on every push and PR"], MED_TEAL),
    ("Native GitHub Integration", ["Results in PRs & Security tab", "Works with GitHub Actions workflows", "Seamless developer experience"], DARK_TEAL),
    ("Flexible Licensing", ["Free for all public repositories", "GHAS license for private/internal repos", "Enterprise-grade security at scale"], MED_TEAL),
]

for i, (title, bullets, color) in enumerate(features):
    left = Inches(0.6) + Inches(i * 3.15)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.4), Inches(2.8), Inches(4.2))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    tb = s.shapes.add_textbox(left + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    bb = s.shapes.add_textbox(left + Inches(0.35), Inches(3.5), Inches(2.1), Inches(2.8))
    add_bullets(bb.text_frame, bullets)

add_notes(s, "GitHub Advanced Security is a suite of security tools designed to help developers find and fix vulnerabilities directly in their development workflow. The key word here is developer-first. These aren\u2019t security tools that slow you down or require you to leave your GitHub environment. GHAS consists of three core components: code scanning, secret scanning, and dependency review. It\u2019s free for public repositories, but for private and internal repositories you\u2019ll need a GHAS license.")
print("Slide 2: What is GHAS")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Core Components
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "GHAS Core Components", "Three pillars of developer-first security")

components = [
    ("Code Scanning", "Find vulnerabilities via static analysis powered by CodeQL.", A_ORANGE,
     ["Static analysis (SAST)", "Powered by CodeQL engine", "SARIF-compatible tools", "Copilot Autofix integration"]),
    ("Secret Scanning", "Detect exposed secrets across your repo, history, issues, and PRs.", A_GREEN,
     ["200+ secret types detected", "150+ partner providers", "Custom regex patterns", "Push protection"]),
    ("Dependency Review", "Understand security impact of dependency changes before merge.", A_BLUE,
     ["PR-level dependency analysis", "Known vulnerability detection", "GitHub Actions integration", "Supply chain protection"]),
]

for i, (title, desc, accent, bullets) in enumerate(components):
    left = Inches(0.6) + Inches(i * 4.2)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.5), Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.6), Inches(3.8), Inches(4.4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tb = s.shapes.add_textbox(left + Inches(0.3), Inches(2.85), Inches(3.2), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.font.name = "Segoe UI"
    db = s.shapes.add_textbox(left + Inches(0.3), Inches(3.4), Inches(3.2), Inches(0.8))
    dtf = db.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(13)
    dp.font.color.rgb = GRAY_TEXT
    dp.font.name = "Segoe UI"
    bb = s.shapes.add_textbox(left + Inches(0.3), Inches(4.3), Inches(3.2), Inches(2.5))
    add_bullets(bb.text_frame, bullets, color=DARK_TEXT)

add_notes(s, "GHAS has three core components that work together: Code Scanning for finding vulnerabilities in your source code, Secret Scanning for detecting leaked credentials, and Dependency Review for managing supply chain risks.")
print("Slide 3: Core Components")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Code Scanning Deep Dive
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "Code Scanning", "Find and fix vulnerabilities before deployment")

items = [
    ("Static Analysis", "Examines source code without executing it, catching SQL injection, XSS, buffer overflows."),
    ("Powered by CodeQL", "GitHub\u2019s semantic engine treats code as data for complex vulnerability queries."),
    ("Third-Party Integration", "Supports any SARIF-compatible tool alongside CodeQL."),
    ("Automated Workflows", "Runs on every push/PR. Results surface directly in pull requests."),
    ("Copilot Autofix", "AI-powered fix suggestions that automatically propose remediation."),
]

lb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
ltf = lb.text_frame
ltf.word_wrap = True
for i, (t, d) in enumerate(items):
    if i > 0:
        ltf.add_paragraph().space_before = Pt(4)
    tp = ltf.add_paragraph() if i > 0 else ltf.paragraphs[0]
    r = tp.add_run()
    r.text = "\u25B6  " + t
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = DARK_TEAL
    r.font.name = "Segoe UI"
    dp = ltf.add_paragraph()
    dp.text = "     " + d
    dp.font.size = Pt(13)
    dp.font.color.rgb = GRAY_TEXT
    dp.font.name = "Segoe UI"
    dp.space_after = Pt(6)

sc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.3), Inches(5.3), Inches(4.5))
sc.fill.solid()
sc.fill.fore_color.rgb = DARK_TEAL
sc.line.fill.background()
st = s.shapes.add_textbox(Inches(7.5), Inches(2.5), Inches(4.7), Inches(0.5))
p = st.text_frame.paragraphs[0]
p.text = "Key Highlights"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Segoe UI"
hb = s.shapes.add_textbox(Inches(7.5), Inches(3.2), Inches(4.7), Inches(3.4))
add_bullets(hb.text_frame, [
    "Supports 10+ programming languages",
    "Thousands of built-in CodeQL queries",
    "Community-contributed query packs",
    "SARIF standard for interoperability",
    "GitHub Actions native integration",
    "Results in Security tab & PR checks",
    "Copilot Autofix for AI remediation",
], font_size=15)

add_notes(s, "Code scanning analyzes your source code to find security vulnerabilities and coding errors before you deploy. It uses static analysis. The primary engine is CodeQL, GitHub\u2019s own semantic code analysis engine. You can also integrate third-party SARIF tools. It integrates with Copilot AutoFix for automatic fix suggestions.")
print("Slide 4: Code Scanning")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Secret Scanning
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "Secret Scanning", "Prevent credential leaks before they cause damage")

ss_features = [
    ("Comprehensive Scanning", A_GREEN,
     ["Scans entire repo including Git history", "Covers issues, PRs, and wikis", "Detects 200+ secret types", "Custom regex pattern support"]),
    ("Partner Program", DARK_TEAL,
     ["150+ service provider partners", "Automatic provider notification", "Credential revocation workflows", "Organization-wide visibility"]),
    ("Push Protection", A_ORANGE,
     ["Blocks pushes containing secrets", "Proactive security enforcement", "Prevents secrets from entering repo", "Bypass with audit trail"]),
]

for i, (title, accent, bullets) in enumerate(ss_features):
    left = Inches(0.6) + Inches(i * 4.2)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(3.8), Inches(4.4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    ch = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.5), Inches(3.8), Inches(0.9))
    ch.fill.solid()
    ch.fill.fore_color.rgb = accent
    ch.line.fill.background()
    tb = s.shapes.add_textbox(left + Inches(0.3), Inches(2.6), Inches(3.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    bb = s.shapes.add_textbox(left + Inches(0.3), Inches(3.7), Inches(3.2), Inches(3.0))
    add_bullets(bb.text_frame, bullets, color=DARK_TEXT, prefix="\u2022  ")

add_notes(s, "Secret scanning is crucial for preventing credential leaks. It searches your entire repository including Git history for sensitive information. It detects 200+ secret types and works with 150+ service providers. Push protection actively blocks pushes containing detected secrets \u2014 proactive security that stops problems before they start.")
print("Slide 5: Secret Scanning")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Dependency Review
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "Dependency Review", "Secure your software supply chain")

db = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(12), Inches(1.0))
dtf = db.text_frame
dtf.word_wrap = True
dp = dtf.paragraphs[0]
dp.text = "Understand how dependency changes in pull requests impact your security posture. Catch supply chain risks before they reach production."
dp.font.size = Pt(16)
dp.font.color.rgb = DARK_TEXT
dp.font.name = "Segoe UI"

dr_features = [
    ("Analyzes Dependency Changes in PRs", A_ORANGE,
     "Scans added/updated dependencies and shows known vulnerabilities in new or updated packages."),
    ("Supply Chain Risk Prevention", A_GREEN,
     "Identifies vulnerable packages before they enter your codebase. Integrates with the GitHub Advisory Database."),
    ("Pre-Merge Security Checks", A_BLUE,
     "Integrates with GitHub Actions to fail PRs introducing high severity vulnerabilities. Gate your merges on security."),
]

for i, (title, color, desc) in enumerate(dr_features):
    left = Inches(0.6) + Inches(i * 4.2)
    ib = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.7), Inches(1.0), Inches(1.0))
    ib.fill.solid()
    ib.fill.fore_color.rgb = color
    ib.line.fill.background()
    tb = s.shapes.add_textbox(left, Inches(4.9), Inches(3.8), Inches(0.7))
    ttf = tb.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(18)
    tp.font.bold = True
    tp.font.color.rgb = DARK_TEXT
    tp.font.name = "Segoe UI"
    db2 = s.shapes.add_textbox(left, Inches(5.6), Inches(3.8), Inches(1.5))
    d2tf = db2.text_frame
    d2tf.word_wrap = True
    d2p = d2tf.paragraphs[0]
    d2p.text = desc
    d2p.font.size = Pt(13)
    d2p.font.color.rgb = GRAY_TEXT
    d2p.font.name = "Segoe UI"

add_notes(s, "Dependency review focuses on supply chain security. When someone creates a PR that adds or updates a dependency, it scans those changes and shows known vulnerabilities. It integrates with GitHub Actions to fail PRs that introduce high severity vulnerabilities. It\u2019s much easier to address a vulnerable dependency in a PR than after it\u2019s in production.")
print("Slide 6: Dependency Review")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Developer Workflow Integration (NEW)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "Developer Workflow Integration", "Security that fits naturally into your daily workflow")

quadrants = [
    ("Automatic Scanning", "Code scanning runs on pushes and PRs without manual intervention. Secret scanning runs constantly in the background. Dependency review triggers on PR creation.",
     Inches(1.5), Inches(2.5)),
    ("Results in GitHub Security Tab", "All alerts appear in the Security tab with actionable info \u2014 where the issue is, what the risk is, and how to fix it. No separate dashboards needed.",
     Inches(7.0), Inches(2.5)),
    ("No Context Switching", "Developers never leave GitHub or learn new tools. Everything is integrated into the platform they already use every day.",
     Inches(1.5), Inches(5.0)),
    ("CI/CD Pipeline Compatible", "For GitHub Actions users, customize when scans run. For others, the CodeQL CLI integrates into any CI/CD environment.",
     Inches(7.0), Inches(5.0)),
]

# Center cross lines
for x1, y1, w, h in [(Inches(6.3), Inches(2.8), Inches(0.06), Inches(3.8)),
                      (Inches(3.0), Inches(4.65), Inches(6.7), Inches(0.06))]:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y1, w, h)
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln.line.fill.background()

for t, d, x, y in quadrants:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_TEAL
    card.line.fill.background()
    tb = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(3.9), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = t
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    db2 = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.8), Inches(3.9), Inches(1.2))
    dtf = db2.text_frame
    dtf.word_wrap = True
    dp2 = dtf.paragraphs[0]
    dp2.text = d
    dp2.font.size = Pt(13)
    dp2.font.color.rgb = LIGHT_TEXT
    dp2.font.name = "Segoe UI"

add_notes(s, "GHAS fits seamlessly into your daily workflow. Once enabled, code scanning runs on pushes and PRs without manual intervention. Secret scanning happens constantly in the background. All alerts appear in the Security tab with actionable information. For teams using GitHub Actions, you can customize when scans run. For teams that don\u2019t, the CodeQL CLI integrates into any CI/CD environment.")
print("Slide 7: Workflow Integration")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Why GHAS? (NEW)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "Why GHAS?", "The value proposition for your development team")

why_items = [
    ("Shift Security Left",
     "Find vulnerabilities while actively writing code, not in production or during separate security reviews. Earlier detection = cheaper and faster fixes.",
     DARK_TEAL, Inches(3.0), Inches(2.3)),
    ("Reduce Developer Friction",
     "Everything integrated into GitHub. No new interfaces to learn, no workflow interruptions. Security becomes a natural part of development.",
     MED_TEAL, Inches(7.3), Inches(2.3)),
    ("Comprehensive Coverage",
     "Code vulnerabilities, credential leaks, and supply chain risks \u2014 all covered in one integrated solution. No gaps in your security posture.",
     MED_TEAL, Inches(3.0), Inches(4.8)),
    ("Integrated Solution",
     "A single platform for all security needs. Unified dashboards, consistent alerts, and centralized management across your organization.",
     DARK_TEAL, Inches(7.3), Inches(4.8)),
]

# Diamond background
diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4.8), Inches(2.6), Inches(3.8), Inches(4.2))
diamond.fill.solid()
diamond.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
diamond.line.fill.background()

for t, d, c, x, y in why_items:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = c
    card.line.fill.background()
    tb = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = t
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    db2 = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.2))
    dtf = db2.text_frame
    dtf.word_wrap = True
    dp2 = dtf.paragraphs[0]
    dp2.text = d
    dp2.font.size = Pt(13)
    dp2.font.color.rgb = LIGHT_TEXT
    dp2.font.name = "Segoe UI"

add_notes(s, "GHAS addresses three fundamental problems: shifting security left so you catch vulnerabilities while actively writing code; reducing developer friction because everything is integrated into GitHub; and providing comprehensive coverage across code vulnerabilities, credential leaks, and supply chain risks in one integrated solution.")
print("Slide 8: Why GHAS")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Summary / Closing
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_TEAL
bg.line.fill.background()

tb = s.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = "Security at Every Stage"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.06))
a.fill.solid()
a.fill.fore_color.rgb = A_BLUE
a.line.fill.background()

sb = s.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.3), Inches(3.0))
stf = sb.text_frame
stf.word_wrap = True
summary_items = [
    "Code Scanning \u2014 Find and fix vulnerabilities with CodeQL + Copilot Autofix",
    "Secret Scanning \u2014 Detect and prevent credential leaks with Push Protection",
    "Dependency Review \u2014 Secure your supply chain before merging",
    "Workflow Integration \u2014 Seamless, automated, zero context-switching",
    "Shift Left \u2014 Comprehensive coverage that reduces friction and catches issues early",
]
for i, t in enumerate(summary_items):
    sp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    sp.text = "\u2713  " + t
    sp.font.size = Pt(20)
    sp.font.color.rgb = LIGHT_TEXT
    sp.font.name = "Segoe UI"
    sp.space_after = Pt(16)

tag = s.shapes.add_textbox(Inches(2.0), Inches(5.6), Inches(9.3), Inches(0.8))
tp = tag.text_frame.paragraphs[0]
tp.text = "Shift left. Secure early. Ship with confidence."
tp.font.size = Pt(24)
tp.font.italic = True
tp.font.color.rgb = A_BLUE
tp.font.name = "Segoe UI"
tp.alignment = PP_ALIGN.CENTER

ft = s.shapes.add_textbox(Inches(2.0), Inches(6.4), Inches(9.3), Inches(0.5))
fp = ft.text_frame.paragraphs[0]
fp.text = "Learn more: github.com/features/security"
fp.font.size = Pt(16)
fp.font.color.rgb = RGBColor(0x99, 0xBB, 0xCC)
fp.font.name = "Segoe UI"
fp.alignment = PP_ALIGN.CENTER

add_notes(s, "In summary, GitHub Advanced Security brings together code scanning, secret scanning, and dependency review for comprehensive security throughout your development lifecycle. These tools integrate seamlessly into your existing workflow so security becomes a natural part of how you build software. Shift left. Secure early. Ship with confidence.")
print("Slide 9: Summary")

# Save
out = os.path.join(os.getcwd(), "GitHub_Advanced_Security.pptx")
prs.save(out)
print(f"\nPPTX saved: {out}")
print(f"Total slides: {len(prs.slides)}")
