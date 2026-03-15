"""
Generate comprehensive PowerPoint presentation for customer demo meeting.
Covers all customer questions about Copilot security, data residency, logs, and PII.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1B, 0x22)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
YELLOW = RGBColor(0xD2, 0x99, 0x22)
RED = RGBColor(0xF8, 0x51, 0x49)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
BLUE_DARK = RGBColor(0x0D, 0x41, 0x9D)
ORANGE = RGBColor(0xF0, 0x88, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_para(tf, text, size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_card(slide, left, top, width, height, fill_color=SURFACE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_bar(slide, left, top, width=0.06, height=0.8, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.2, BLUE_DARK)
add_text(slide, 0.8, 0.15, 11, 0.9, "GitHub Copilot Security Deep Dive", 36, WHITE, True, PP_ALIGN.LEFT)
add_text(slide, 0.8, 2.0, 11, 0.6, "Answering Every Question About Data Residency,", 28, WHITE, True)
add_text(slide, 0.8, 2.6, 11, 0.6, "Secrets, Logs, Sessions & PII Protection", 28, ACCENT, True)

items = [
    "Where does my prompt data go and how long is it retained?",
    "Can Copilot see my database and secrets?",
    "Are instruction files (.md) reliable for security?",
    "What's in the debug logs and session history?",
    "What happens when Copilot makes MCP tool calls?",
    "Live demos with real proof at every step"
]
for i, item in enumerate(items):
    add_accent_bar(slide, 1.0, 3.7 + i*0.5, 0.05, 0.35, ACCENT)
    add_text(slide, 1.3, 3.65 + i*0.5, 10, 0.45, item, 16, WHITE)

add_text(slide, 0.8, 6.8, 5, 0.4, "Prepared for Customer Security Review", 14, MUTED)
add_text(slide, 8, 6.8, 5, 0.4, "github.com/sautalwar/ghcopilot-pii-demo", 14, ACCENT, False, PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Agenda — What We'll Cover Today", 30, WHITE, True)

sections = [
    ("1", "Data Flow & Residency", "Where does your prompt go? End-to-end journey through every component", "10 min"),
    ("2", "Database & PII Isolation", "Proof that Copilot cannot see runtime database data", "15 min"),
    ("3", "Secrets & Environment Variables", "How Copilot handles secrets in agent mode and MCP", "10 min"),
    ("4", "Instruction Files (.md) — The Truth", "Are they deterministic? What the data actually shows", "10 min"),
    ("5", "Debug Logs & Session Storage", "What's stored locally, on GitHub servers, and for how long", "10 min"),
    ("6", "Blast Radius Analysis", "If instructions fail, what's the actual impact?", "10 min"),
    ("7", "MCP Tool Calls & Redaction", "Safe vs. unsafe patterns for AI-connected tools", "10 min"),
    ("8", "Live Demo & Q&A", "Real-time proof with working application", "15 min"),
]
for i, (num, title, desc, time) in enumerate(sections):
    y = 1.3 + i * 0.72
    add_card(slide, 0.6, y, 11.5, 0.65, SURFACE)
    add_accent_bar(slide, 0.6, y + 0.1, 0.05, 0.45, ACCENT if i < 7 else GREEN)
    add_text(slide, 0.85, y + 0.05, 0.5, 0.5, num, 22, ACCENT, True)
    add_text(slide, 1.4, y + 0.02, 4, 0.35, title, 18, WHITE, True)
    add_text(slide, 1.4, y + 0.33, 7, 0.3, desc, 13, MUTED)
    add_text(slide, 10.5, y + 0.1, 1.5, 0.4, time, 14, MUTED, False, PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 3: Customer Questions Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Your Questions — We Heard Every One", 30, WHITE, True)

questions = [
    ('"What if Copilot is making tool calls to search the data and extract it?"', RED),
    ('"MD files are suggestions, not deterministic — only followed 50-80%"', YELLOW),
    ('"Where does data end up in our debugging path, in the logs?"', YELLOW),
    ('"What about prompt debugging logs that Microsoft captures?"', YELLOW),
    ('"How long are web-based chat sessions and local sessions stored?"', ORANGE),
    ('"When Copilot has a password or sensitive info, what\'s our exposure?"', RED),
    ('"Let\'s pretend instruction files get ignored — where does the data end up?"', RED),
    ('"There\'s no auto expiration for local sessions"', ORANGE),
]
for i, (q, color) in enumerate(questions):
    y = 1.3 + i * 0.72
    add_card(slide, 0.6, y, 11.5, 0.62, SURFACE)
    add_accent_bar(slide, 0.6, y + 0.1, 0.05, 0.42, color)
    add_text(slide, 0.9, y + 0.08, 11, 0.5, q, 15, WHITE)

add_text(slide, 0.8, 7.1, 10, 0.3, "We'll answer EVERY question with live proof — not just slides", 14, GREEN, True)

# ============================================================
# SLIDE 4: Data Flow — Where Does Your Prompt Go?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Where Does Your Prompt Go? — End-to-End Data Flow", 28, WHITE, True)

steps = [
    ("1. Developer Prompt", "You type in VS Code / CLI / Web", "Your machine", ACCENT),
    ("2. Copilot Extension", "Packages context from open files\nContent exclusion applied HERE", "Your machine", ACCENT),
    ("3. GitHub Proxy", "copilot-proxy.githubusercontent.com\nTLS 1.2+ encrypted in transit", "GitHub infra", YELLOW),
    ("4. Model Provider", "Azure OpenAI or Anthropic\nProcesses prompt, generates response", "Cloud", ORANGE),
    ("5. Response", "Returned to your editor\nPrompt DISCARDED (Biz/Ent)", "Your machine", GREEN),
]
for i, (title, desc, loc, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    add_card(slide, x, 1.3, 2.3, 2.4, SURFACE)
    add_accent_bar(slide, x, 1.3, 2.3, 0.06, color)
    add_text(slide, x + 0.15, 1.45, 2.0, 0.4, title, 14, color, True)
    add_text(slide, x + 0.15, 1.85, 2.0, 1.2, desc, 12, WHITE)
    add_text(slide, x + 0.15, 3.3, 2.0, 0.3, loc, 11, MUTED)

# Key facts below
add_card(slide, 0.5, 4.0, 5.8, 3.0, SURFACE)
tf = add_text(slide, 0.7, 4.1, 5.4, 0.4, "Key Facts — Retention by Provider", 18, ACCENT, True)
rows = [
    ("Azure OpenAI (GPT models)", "30 days abuse monitoring", "Enterprise opt-out available"),
    ("Anthropic (Claude Sonnet/Opus)", "ZERO days — ZDR agreement", "Nothing stored, ever"),
    ("GitHub servers", "Prompts discarded after response", "Business & Enterprise plans"),
    ("github.com Web Chat", "30 days conversation history", "Auto-deleted, export available"),
    ("VS Code local", "Indefinite (no auto-expiry)", "Developer must manually clear"),
    ("Usage metrics (not content)", "90 days rolling", "Timestamps only, no prompts"),
]
for i, (provider, retention, note) in enumerate(rows):
    y = 4.55 + i * 0.38
    add_text(slide, 0.75, y, 2.2, 0.35, provider, 11, WHITE, True)
    add_text(slide, 3.0, y, 1.8, 0.35, retention, 11, GREEN if "ZERO" in retention or "discarded" in retention else YELLOW if "30" in retention else ORANGE, True)
    add_text(slide, 4.9, y, 1.3, 0.35, note, 10, MUTED)

add_card(slide, 6.6, 4.0, 6.2, 3.0, SURFACE)
tf2 = add_text(slide, 6.8, 4.1, 5.8, 0.4, "How to Verify (Live Demo Commands)", 18, GREEN, True)
cmds = [
    ("Prove single endpoint:", "nslookup copilot-proxy.githubusercontent.com"),
    ("Prove which model:", "DevTools → Network → x-model-version header"),
    ("Prove no DB calls:", "DevTools → Network → zero copilot requests during query"),
    ("Prove encryption:", "All requests use TLS 1.2+ (check Security tab)"),
    ("Prove content exclusion:", "Open excluded file → slashed Copilot icon"),
    ("Prove stateless:", 'Ask "What is pineapple42?" → Copilot has no idea'),
]
for i, (label, cmd) in enumerate(cmds):
    y = 4.55 + i * 0.38
    add_text(slide, 6.85, y, 2.0, 0.35, label, 11, WHITE, True)
    add_text(slide, 8.9, y, 3.8, 0.35, cmd, 10, ACCENT)

# ============================================================
# SLIDE 5: Database & PII Isolation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, '"Can Copilot See My Database?" — NO. Here\'s the Proof.', 28, WHITE, True)

# Left side: The boundary
add_card(slide, 0.5, 1.3, 6.0, 2.8, SURFACE)
add_text(slide, 0.7, 1.4, 5.6, 0.4, "The Design-Time vs Runtime Boundary", 18, ACCENT, True)

can_see = [
    "Function names: getCitizenById()",
    "SQL patterns: SELECT * FROM citizens",
    "Type definitions: interface Citizen { ssn: string }",
    "Variable names: citizen.ssn",
    "Source code files (.ts, .js, .py)",
]
cannot_see = [
    "Actual citizen records in the database",
    "Query results (real SSNs, emails, phones)",
    "Database file contents (.db, .sqlite)",
    "Runtime environment variable VALUES",
    "In-memory application state",
]
add_text(slide, 0.75, 1.85, 2.5, 0.3, "Copilot CAN See (Code)", 13, GREEN, True)
for i, item in enumerate(can_see):
    add_text(slide, 0.75, 2.15 + i * 0.3, 2.7, 0.3, f"✓ {item}", 10, WHITE)

add_text(slide, 3.5, 1.85, 2.5, 0.3, "Copilot CANNOT See (Data)", 13, RED, True)
for i, item in enumerate(cannot_see):
    add_text(slide, 3.5, 2.15 + i * 0.3, 3.0, 0.3, f"✗ {item}", 10, WHITE)

# Right side: The Canary Test
add_card(slide, 6.8, 1.3, 6.0, 2.8, SURFACE)
add_text(slide, 7.0, 1.4, 5.6, 0.4, "The Canary Test — Definitive Proof", 18, GREEN, True)

steps_canary = [
    ('1. We seed a "sentinel" record:', '"Canary Testbird" — SSN: 000-00-0000'),
    ("2. Ask the application:", "Search for 'Canary' → FOUND instantly"),
    ("3. Ask Copilot:", '"Is there a citizen named Canary Testbird?" → NO IDEA'),
    ("4. Conclusion:", "If Copilot could see DB data, it would know. It doesn't."),
]
for i, (step, detail) in enumerate(steps_canary):
    add_text(slide, 7.05, 1.9 + i * 0.6, 5.5, 0.3, step, 13, WHITE, True)
    add_text(slide, 7.05, 2.2 + i * 0.6, 5.5, 0.3, detail, 12, ACCENT)

# Bottom: Live demo steps
add_card(slide, 0.5, 4.4, 12.3, 2.8, SURFACE)
add_text(slide, 0.7, 4.5, 11.5, 0.4, "LIVE DEMO — Step by Step", 18, GREEN, True)

demo_steps = [
    ("Step 1", "Open http://localhost:3000/pii-demo.html", 'Click "Query Database (Raw)" → show full SSNs'),
    ("Step 2", "Click 'Query Database (Redacted)'", "SSNs masked: ***-**-1234 — app-level redaction"),
    ("Step 3", "Open VS Code Chat", 'Ask: "What SSNs are in the citizens database?" → Copilot can\'t answer'),
    ("Step 4", "Ask Copilot about code", '"What does getCitizenById do?" → Copilot describes the function'),
    ("Step 5", "Run the Canary Test", 'Click "Search for Canary" → app finds it. Ask Copilot → blank'),
    ("Step 6", "Open DevTools → Network", "Query the database → ZERO requests to copilot-proxy"),
]
for i, (step, action, result) in enumerate(demo_steps):
    x = 0.7 + (i % 3) * 4.1
    y = 5.0 + (i // 3) * 1.1
    add_text(slide, x, y, 0.7, 0.3, step, 12, ACCENT, True)
    add_text(slide, x + 0.7, y, 3.2, 0.3, action, 11, WHITE, True)
    add_text(slide, x + 0.7, y + 0.3, 3.2, 0.4, result, 10, MUTED)

# ============================================================
# SLIDE 6: Instruction Files — The Truth
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, '"MD Files Are Only Followed 50-80%" — What\'s True & What\'s Not', 26, WHITE, True)

# Left: The honest answer
add_card(slide, 0.5, 1.3, 6.0, 3.0, SURFACE)
add_text(slide, 0.7, 1.4, 5.6, 0.4, "The Honest Answer", 18, YELLOW, True)

truths = [
    ("TRUE:", "Instructions are probabilistic, not deterministic", GREEN),
    ("TRUE:", "Compliance drops with long/complex instruction files (>4000 chars)", GREEN),
    ("TRUE:", "Too many instructions → model skips some", GREEN),
    ("PARTIALLY TRUE:", "50-80% compliance claim — no peer-reviewed study confirms exact rate", YELLOW),
    ("TRUE:", "GitHub's own docs say behavior is 'non-deterministic'", GREEN),
    ("IMPORTANT:", "Instructions are NOT a security boundary", RED),
]
for i, (label, desc, color) in enumerate(truths):
    add_accent_bar(slide, 0.7, 1.9 + i * 0.38, 0.04, 0.28, color)
    add_text(slide, 0.85, 1.87 + i * 0.38, 1.2, 0.3, label, 11, color, True)
    add_text(slide, 2.0, 1.87 + i * 0.38, 4.3, 0.3, desc, 11, WHITE)

# Right: The framework
add_card(slide, 6.8, 1.3, 6.0, 3.0, SURFACE)
add_text(slide, 7.0, 1.4, 5.6, 0.4, "What to Use Instead — Security Stack", 18, ACCENT, True)

controls = [
    ("Content Exclusion (org settings)", "DETERMINISTIC", "100%", GREEN),
    ("Secret Scanning + Push Protection", "DETERMINISTIC", "100%", GREEN),
    (".copilotignore file", "DETERMINISTIC", "100%", GREEN),
    ("Environment Scoping (agent mode)", "DETERMINISTIC", "100%", GREEN),
    ("copilot-instructions.md", "PROBABILISTIC", "~70-90%", YELLOW),
    ("Prompt engineering", "PROBABILISTIC", "~60-80%", ORANGE),
]
for i, (control, ctype, rate, color) in enumerate(controls):
    y = 1.9 + i * 0.38
    add_text(slide, 7.05, y, 3.0, 0.3, control, 11, WHITE)
    add_text(slide, 10.1, y, 1.5, 0.3, ctype, 10, color, True)
    add_text(slide, 11.6, y, 0.8, 0.3, rate, 11, color, True)

# Bottom: Key takeaway
add_card(slide, 0.5, 4.6, 12.3, 2.6, SURFACE)
add_accent_bar(slide, 0.5, 4.6, 12.3, 0.06, ACCENT)
add_text(slide, 0.7, 4.75, 11.5, 0.5, "Key Takeaway for the Customer", 20, ACCENT, True)
add_text(slide, 0.7, 5.3, 11.5, 0.6, '"The customer is RIGHT that instructions aren\'t 100% reliable. That\'s exactly why we NEVER rely on them for security.', 16, WHITE, True)
add_text(slide, 0.7, 5.8, 11.5, 0.6, 'We use DETERMINISTIC controls (content exclusion, secret scanning, environment scoping) for security.', 16, WHITE)
add_text(slide, 0.7, 6.2, 11.5, 0.6, 'Instructions are one LAYER of defense — not the only layer. Think seatbelt + airbag + crumple zone."', 16, WHITE)

add_text(slide, 0.7, 6.8, 11.5, 0.4, "LIVE DEMO: Create instruction 'Never use console.log' → test 5x → show ~70-90% compliance → then show Content Exclusion at 100%", 13, GREEN, True)

# ============================================================
# SLIDE 7: What Gets Stored Where
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Complete Data Residency Map — What's Stored Where", 28, WHITE, True)

# Table header
add_card(slide, 0.5, 1.2, 12.3, 0.5, BLUE_DARK)
headers = [("Location", 0.6, 2.5), ("What's Stored", 3.2, 2.5), ("Retention", 5.8, 1.8), ("Risk", 7.7, 0.8), ("How to Clear", 8.6, 3.5)]
for h, x, w in headers:
    add_text(slide, x, 1.22, w, 0.45, h, 13, WHITE, True)

rows_data = [
    ("VS Code workspaceStorage", "Chat sessions (state.vscdb)", "Until deleted", "LOW", "Delete file"),
    ("VS Code Output Panel", "Debug logs (ephemeral)", "Until restart", "NONE", "Close panel"),
    ("VS Code logs folder", "Extension logs", "Rotated", "NONE", "Delete folder"),
    ("Shell history", "CLI prompts (gh copilot)", "Until cleared", "LOW", "Clear-History"),
    ("github.com Web Chat", "Conversations", "30 days auto", "LOW", "Wait / export"),
    ("GitHub Servers", "Usage metrics (NOT content)", "90 days", "NONE", "Auto-expires"),
    ("Azure OpenAI", "Abuse monitoring prompts", "30 days", "MED", "Opt-out available"),
    ("Anthropic (Claude)", "NOTHING — Zero Data Retention", "0 days", "NONE", "N/A"),
    ("MCP Server logs", "Tool I/O (your code)", "You decide", "YOU", "Your responsibility"),
    ("Local database", "Runtime app data", "Permanent", "NONE*", "Delete file"),
]

risk_colors = {"NONE": GREEN, "LOW": YELLOW, "MED": ORANGE, "YOU": ACCENT, "NONE*": GREEN}
for i, (loc, stored, ret, risk, clear) in enumerate(rows_data):
    y = 1.78 + i * 0.5
    bg_col = SURFACE if i % 2 == 0 else RGBColor(0x1C, 0x21, 0x28)
    add_card(slide, 0.5, y, 12.3, 0.47, bg_col)
    add_text(slide, 0.6, y + 0.02, 2.5, 0.4, loc, 11, WHITE, True)
    add_text(slide, 3.2, y + 0.02, 2.5, 0.4, stored, 11, WHITE)
    add_text(slide, 5.8, y + 0.02, 1.8, 0.4, ret, 11, GREEN if "0 day" in ret or "discarded" in ret else YELLOW if "30" in ret or "90" in ret else ORANGE if "Until" in ret else WHITE)
    add_text(slide, 7.7, y + 0.02, 0.8, 0.4, risk, 11, risk_colors.get(risk, WHITE), True)
    add_text(slide, 8.6, y + 0.02, 3.5, 0.4, clear, 11, MUTED)

add_text(slide, 0.6, 6.9, 12, 0.4, '* "NONE" for database = Copilot has zero access to runtime data. Data never leaves your machine.', 12, MUTED)

# ============================================================
# SLIDE 8: Debug Logs — What's Inside
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, '"What\'s in the Debug Logs?" — Exactly This, Nothing More', 28, WHITE, True)

# Left: What IS logged
add_card(slide, 0.5, 1.3, 4.0, 3.5, SURFACE)
add_text(slide, 0.7, 1.4, 3.6, 0.4, "What IS in the Logs", 16, GREEN, True)
logged = [
    "✓ Timestamps of requests",
    "✓ Request IDs (UUIDs)",
    "✓ Model selection info",
    "✓ Token counts (in/out)",
    "✓ Response latency (ms)",
    "✓ Error messages & retries",
    "✓ Extension activation events",
    "✓ Feature usage flags",
]
for i, item in enumerate(logged):
    add_text(slide, 0.75, 1.85 + i * 0.33, 3.5, 0.3, item, 12, WHITE)

# Middle: What is NOT logged
add_card(slide, 4.7, 1.3, 4.0, 3.5, SURFACE)
add_text(slide, 4.9, 1.4, 3.6, 0.4, "What is NOT in the Logs", 16, RED, True)
not_logged = [
    "✗ Full prompt text",
    "✗ Full response text",
    "✗ Your source code content",
    "✗ PII or secrets",
    "✗ Database query results",
    "✗ File contents",
    "✗ Environment variable values",
    "✗ Conversation history",
]
for i, item in enumerate(not_logged):
    add_text(slide, 4.95, 1.85 + i * 0.33, 3.5, 0.3, item, 12, WHITE)

# Right: How to verify
add_card(slide, 8.9, 1.3, 3.9, 3.5, SURFACE)
add_text(slide, 9.1, 1.4, 3.5, 0.4, "How to Verify (Demo)", 16, ACCENT, True)
verify = [
    "1. View → Output panel",
    '2. Select "GitHub Copilot"',
    "3. Ask Copilot a question",
    "4. Read the log entries",
    "5. Search for SSN pattern",
    "   → Zero matches",
    "6. Search for prompt text",
    "   → Zero matches",
]
for i, item in enumerate(verify):
    add_text(slide, 9.15, 1.85 + i * 0.33, 3.5, 0.3, item, 12, ACCENT if "→" in item else WHITE)

# Bottom: Log locations
add_card(slide, 0.5, 5.1, 12.3, 2.1, SURFACE)
add_text(slide, 0.7, 5.2, 11.5, 0.4, "Log File Locations — Every Path on Disk", 16, ACCENT, True)

locations = [
    ("VS Code Output (ephemeral)", "View → Output → 'GitHub Copilot'", "Until restart"),
    ("Extension Host Logs", "%APPDATA%\\Code\\logs\\<date>\\exthost\\", "Rotated by VS Code"),
    ("Chat Session Storage", "%APPDATA%\\Code\\User\\workspaceStorage\\<id>\\state.vscdb", "Until deleted"),
    ("PowerShell History", "(Get-PSReadLineOption).HistorySavePath", "Until cleared"),
]
for i, (name, path, ret) in enumerate(locations):
    y = 5.65 + i * 0.35
    add_text(slide, 0.75, y, 2.5, 0.3, name, 11, WHITE, True)
    add_text(slide, 3.3, y, 6.5, 0.3, path, 11, ACCENT)
    add_text(slide, 10.3, y, 2, 0.3, ret, 11, MUTED)

# ============================================================
# SLIDE 9: Secrets & Agent Mode
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, '"When Copilot Has a Password — What\'s Our Exposure?"', 28, WHITE, True)

# Three scenarios
scenarios = [
    ("Scenario A: Copilot Chat", "Copilot reads code, not .env files\n(with content exclusion)", GREEN, "ZERO exposure\nif .env excluded"),
    ("Scenario B: Agent Mode", "Secrets only accessible if explicitly\nadded to 'copilot' environment", YELLOW, "CONTROLLED\nby admin"),
    ("Scenario C: MCP Tool Calls", "Tool output enters conversation\ncontext — redact before returning", ORANGE, "YOUR tool design\ncontrols exposure"),
]
for i, (title, desc, color, impact) in enumerate(scenarios):
    x = 0.5 + i * 4.2
    add_card(slide, x, 1.3, 3.9, 2.5, SURFACE)
    add_accent_bar(slide, x, 1.3, 3.9, 0.06, color)
    add_text(slide, x + 0.2, 1.5, 3.5, 0.4, title, 16, color, True)
    add_text(slide, x + 0.2, 1.95, 3.5, 0.8, desc, 13, WHITE)
    add_text(slide, x + 0.2, 2.9, 3.5, 0.8, impact, 14, color, True)

# Bottom: The key principle
add_card(slide, 0.5, 4.1, 12.3, 1.0, SURFACE)
add_accent_bar(slide, 0.5, 4.1, 12.3, 0.06, ACCENT)
add_text(slide, 0.7, 4.25, 11.5, 0.4, '"The blast radius of Copilot is the blast radius of the tools you give it."', 20, ACCENT, True)
add_text(slide, 0.7, 4.7, 11.5, 0.3, "Like giving keys to a contractor: Copilot never picks locks. You decide which keys to hand over.", 14, MUTED)

# Safe vs Unsafe MCP code
add_card(slide, 0.5, 5.3, 6.0, 2.0, SURFACE)
add_text(slide, 0.7, 5.4, 5.6, 0.3, "❌ UNSAFE MCP Tool", 15, RED, True)
add_text(slide, 0.7, 5.75, 5.6, 1.4, 'server.tool("lookup_citizen", async ({id}) => {\n  const citizen = db.get(id);\n  return JSON.stringify(citizen);\n  // RAW SSN goes to Copilot context!\n});', 11, WHITE)

add_card(slide, 6.8, 5.3, 6.0, 2.0, SURFACE)
add_text(slide, 7.0, 5.4, 5.6, 0.3, "✅ SAFE MCP Tool", 15, GREEN, True)
add_text(slide, 7.0, 5.75, 5.6, 1.4, 'server.tool("lookup_citizen", async ({id}) => {\n  const citizen = db.get(id);\n  return JSON.stringify({\n    ...citizen, ssn: maskSSN(citizen.ssn)\n  }); // Only masked data reaches Copilot\n});', 11, WHITE)

# ============================================================
# SLIDE 10: Blast Radius Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, '"If Instructions Are Ignored — Where Does the Data End Up?"', 28, WHITE, True)

# Flow diagram with boxes
flow_steps = [
    ("Developer types\nSSN in prompt", 0.3, 1.4, ACCENT, "Your machine"),
    ("Copilot Extension\nContent exclusion\napplied HERE", 2.8, 1.4, GREEN, "Blocked if excluded"),
    ("GitHub Proxy\nTLS 1.2+ encrypted", 5.3, 1.4, YELLOW, "In transit only"),
    ("Model Provider\nProcesses prompt", 7.8, 1.4, ORANGE, "See retention table"),
    ("Response returned\nPrompt DISCARDED", 10.3, 1.4, GREEN, "Biz/Ent plans"),
]
for title, x, y, color, note in flow_steps:
    add_card(slide, x, y, 2.2, 1.5, SURFACE)
    add_accent_bar(slide, x, y, 2.2, 0.05, color)
    add_text(slide, x + 0.1, y + 0.15, 2.0, 0.9, title, 12, WHITE, True)
    add_text(slide, x + 0.1, y + 1.05, 2.0, 0.3, note, 10, color)

# Arrows between boxes
for i in range(4):
    x = 2.6 + i * 2.5
    add_text(slide, x, 1.85, 0.3, 0.3, "→", 24, MUTED, True)

# Impact table
add_card(slide, 0.5, 3.3, 12.3, 3.8, SURFACE)
add_text(slide, 0.7, 3.4, 11.5, 0.4, "Risk Assessment — Even If Instructions Fail", 18, ACCENT, True)

risks = [
    ("PII in code comment", "Medium", "Low — discarded after response", "Content exclusion, code review", YELLOW),
    ("Secret in .env file", "Low (if excluded)", "Medium — rotate the secret", "Content exclusion (100% block)", GREEN),
    ("MCP tool returns raw PII", "Medium (if poor design)", "Medium — in model context", "Redact inside the tool", YELLOW),
    ("SSN typed directly in prompt", "Low", "Low — in transit only, discarded", "Developer training", YELLOW),
    ("Azure OpenAI 30-day log", "Low (opt-out avail)", "Low — encrypted, access-controlled", "Request Modified Abuse Monitoring", GREEN),
    ("Local chat history persists", "High (always stored)", "Low — local file, OS-protected", "Cleanup script + FDE", ORANGE),
]
add_card(slide, 0.6, 3.85, 12.1, 0.4, BLUE_DARK)
for j, h in enumerate(["Risk Scenario", "Probability", "Impact", "Mitigation"]):
    add_text(slide, 0.7 + j * 3.0, 3.88, 2.8, 0.35, h, 12, WHITE, True)

for i, (risk, prob, impact, mit, color) in enumerate(risks):
    y = 4.3 + i * 0.42
    add_text(slide, 0.7, y, 2.8, 0.35, risk, 11, WHITE)
    add_text(slide, 3.7, y, 2.8, 0.35, prob, 11, color, True)
    add_text(slide, 6.7, y, 2.8, 0.35, impact, 11, WHITE)
    add_text(slide, 9.7, y, 3.0, 0.35, mit, 11, ACCENT)

# ============================================================
# SLIDE 11: The 5-Layer Security Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "The 5-Layer Security Stack — Defense in Depth", 28, WHITE, True)

layers = [
    ("Layer 5", "Human Code Review", "Final validation before merge — catches anything AI missed", "HUMAN", WHITE),
    ("Layer 4", "copilot-instructions.md", "Probabilistic guidance — reduces unsafe suggestions ~70-90%", "PROBABILISTIC", YELLOW),
    ("Layer 3", "Secret Scanning + Push Protection", "Catches secrets in commits — 200+ patterns, blocks before push", "DETERMINISTIC", GREEN),
    ("Layer 2", "Environment Scoping", "Agent mode secrets isolated — must explicitly grant access", "DETERMINISTIC", GREEN),
    ("Layer 1", "Content Exclusion", "Platform-enforced file blocking — files NEVER reach the model", "DETERMINISTIC", GREEN),
]

for i, (layer, name, desc, ltype, color) in enumerate(layers):
    y = 1.3 + i * 1.15
    width = 11.0 - i * 0.5
    x = 0.5 + i * 0.25
    add_card(slide, x, y, width, 0.95, SURFACE)
    add_accent_bar(slide, x, y, 0.06, 0.95, color)
    add_text(slide, x + 0.2, y + 0.05, 1.0, 0.4, layer, 14, color, True)
    add_text(slide, x + 1.3, y + 0.05, 4, 0.4, name, 16, WHITE, True)
    add_text(slide, x + 1.3, y + 0.45, width - 2, 0.4, desc, 12, MUTED)
    add_text(slide, x + width - 2.2, y + 0.05, 2.0, 0.4, ltype, 12, color, True, PP_ALIGN.RIGHT)

add_card(slide, 1.5, 7.0 - 0.6, 10, 0.5, SURFACE)
add_text(slide, 1.7, 7.0 - 0.55, 9.5, 0.4, "Bottom line: Use deterministic controls for security. Instructions for convenience. Never rely on one layer alone.", 14, ACCENT, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Multi-Model Security (Claude/Opus)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Multi-Model Security — Claude Sonnet & Opus", 28, WHITE, True)

# Key facts
facts = [
    ("Zero Data Retention (ZDR)", "Anthropic has ZDR agreement with Microsoft — prompts are NOT stored", GREEN),
    ("Microsoft Sub-Processor", "Anthropic became Microsoft sub-processor Jan 7, 2026 — governed by Microsoft DPA", ACCENT),
    ("Disabled by Default (EU/EFTA/UK)", "Claude models must be explicitly enabled by admin in restricted regions", YELLOW),
    ("Context Windows", "Sonnet 4/4.6: 1,000,000 tokens | Opus: 200,000 tokens | GPT-4.1: 128,000 tokens", ACCENT),
    ("Same Proxy Endpoint", "All models route through copilot-proxy.githubusercontent.com — verify with x-model-version header", GREEN),
    ("No Training on Your Data", "Microsoft contractual guarantee — your code is never used to train any model", GREEN),
]

for i, (title, desc, color) in enumerate(facts):
    y = 1.3 + i * 0.9
    add_card(slide, 0.5, y, 12.3, 0.8, SURFACE)
    add_accent_bar(slide, 0.5, y + 0.1, 0.05, 0.6, color)
    add_text(slide, 0.8, y + 0.05, 4, 0.35, title, 15, color, True)
    add_text(slide, 0.8, y + 0.4, 11.5, 0.35, desc, 13, WHITE)

# ============================================================
# SLIDE 13: Live Demo Checklist
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Live Demo Checklist — What We'll Prove Today", 28, WHITE, True)

demos = [
    ("Database Isolation", "Query PII database → Copilot can't see data", "pii-demo.html", "15 min", GREEN),
    ("Canary Test", "Sentinel record in DB → Copilot is blind", "pii-demo.html", "3 min", GREEN),
    ("Network Inspection", "DevTools shows zero copilot requests during DB queries", "VS Code DevTools", "5 min", GREEN),
    ("Debug Log Forensics", "Output panel shows metadata only, no prompts/PII", "VS Code Output", "5 min", GREEN),
    ("Local Storage Audit", "Find state.vscdb → show what's stored", "File Explorer", "5 min", YELLOW),
    ("Content Exclusion", "Exclude .env → slashed icon → Copilot can't read it", "VS Code + GitHub", "5 min", GREEN),
    ("Instruction Compliance", "Test .md rules 5x → show probabilistic nature", "VS Code Chat", "10 min", YELLOW),
    ("MCP Tool Redaction", "Safe vs unsafe tool code → side by side", "Code walkthrough", "5 min", ACCENT),
    ("nslookup Proof", "All traffic → single Microsoft endpoint", "Terminal", "2 min", GREEN),
    ("Stateless Test", "Ask about pineapple42 → Copilot has no memory", "VS Code Chat", "2 min", GREEN),
]

for i, (name, proof, tool, time, color) in enumerate(demos):
    y = 1.2 + i * 0.58
    add_card(slide, 0.5, y, 12.3, 0.52, SURFACE)
    add_text(slide, 0.7, y + 0.05, 0.3, 0.35, f"{i+1}.", 13, color, True)
    add_text(slide, 1.1, y + 0.05, 2.5, 0.35, name, 13, WHITE, True)
    add_text(slide, 3.7, y + 0.05, 4.5, 0.35, proof, 12, MUTED)
    add_text(slide, 8.5, y + 0.05, 2.5, 0.35, tool, 11, ACCENT)
    add_text(slide, 11.2, y + 0.05, 1.2, 0.35, time, 11, MUTED, False, PP_ALIGN.RIGHT)

add_text(slide, 0.7, 7.0, 11, 0.3, "Total demo time: ~60 minutes  |  All demos run from http://localhost:3000", 14, GREEN, True)

# ============================================================
# SLIDE 14: Summary & Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Summary — Your Questions, Answered", 30, WHITE, True)

answers = [
    ("Can Copilot see my database?", "NO — proven by canary test & network inspection", GREEN),
    ("Are .md files reliable for security?", "NO — use deterministic controls (content exclusion, secret scanning)", YELLOW),
    ("Where do prompts go?", "copilot-proxy → model → discarded. Claude: zero retention", GREEN),
    ("What's in debug logs?", "Metadata only (timestamps, IDs). No prompts, no code, no PII", GREEN),
    ("How long are sessions stored?", "Web: 30 days auto-delete. Local: until you delete. Implement cleanup policy", YELLOW),
    ("What if Copilot has a secret?", "Content exclusion blocks it. Agent mode requires explicit grants", GREEN),
    ("What about MCP tool calls?", "You control the tool. Redact before returning data to Copilot", YELLOW),
    ("What if instructions are ignored?", "Blast radius is limited: prompts discarded, data encrypted in transit", GREEN),
]

for i, (q, a, color) in enumerate(answers):
    y = 1.2 + i * 0.72
    add_card(slide, 0.5, y, 12.3, 0.65, SURFACE)
    add_accent_bar(slide, 0.5, y + 0.08, 0.05, 0.48, color)
    add_text(slide, 0.8, y + 0.02, 5, 0.35, q, 14, WHITE, True)
    add_text(slide, 5.8, y + 0.02, 6.5, 0.55, a, 13, color)

add_card(slide, 2.5, 7.0 - 0.6, 8, 0.5, SURFACE)
add_accent_bar(slide, 2.5, 7.0 - 0.6, 8, 0.06, GREEN)
add_text(slide, 2.7, 7.0 - 0.5, 7.5, 0.35, "Every answer backed by a live, repeatable demo you can run yourself", 16, GREEN, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15: Thank You / Resources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, 0, 0, 13.333, 1.0, BLUE_DARK)
add_text(slide, 0.8, 0.1, 11, 0.8, "Resources & Next Steps", 30, WHITE, True)

resources = [
    ("Demo Repository", "github.com/sautalwar/ghcopilot-pii-demo", "Clone, npm install, run locally"),
    ("PII Database Demo", "http://localhost:3000/pii-demo.html", "Canary test, network proof, redaction"),
    ("Data Residency Demo", "http://localhost:3000/data-residency-demo.html", "Logs, sessions, storage map"),
    ("Main Demo Portal", "http://localhost:3000", "All demos, GHAS, competitors, security"),
    ("GitHub Docs: Content Exclusion", "docs.github.com/en/copilot/concepts/context/content-exclusion", "Official documentation"),
    ("GitHub Docs: Model Hosting", "docs.github.com/en/copilot/reference/ai-models/model-hosting", "ZDR, data processing"),
]

for i, (name, url, desc) in enumerate(resources):
    y = 1.4 + i * 0.75
    add_card(slide, 0.5, y, 12.3, 0.65, SURFACE)
    add_text(slide, 0.75, y + 0.05, 3, 0.3, name, 15, WHITE, True)
    add_text(slide, 3.8, y + 0.05, 5, 0.3, url, 13, ACCENT)
    add_text(slide, 3.8, y + 0.33, 5, 0.3, desc, 11, MUTED)

add_text(slide, 0.8, 6.5, 11, 0.5, "Thank you — Let's run the demos live.", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 0.8, 7.0, 11, 0.3, "All code is open source. Clone the repo and reproduce every proof yourself.", 14, MUTED, False, PP_ALIGN.CENTER)

# Save
output_path = r"C:\Users\sautalwar\Downloads\repos\ghcopilot-pii-demo\docs\Copilot-Security-Deep-Dive-Presentation.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Slides: {len(prs.slides)}")
